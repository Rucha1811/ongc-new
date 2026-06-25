from fastapi import APIRouter, Depends, HTTPException, Form
from fastapi.responses import Response
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select
from sqlalchemy.orm import selectinload
from app.database import get_db
from app.models.base import ReportTemplate, ReportPeriod, ReportSubmission, User
from app.auth.deps import get_current_user
from datetime import datetime, date
import json

router = APIRouter()

# ─── Helper ───
def admin_or_ops(user):
    role = user.role.name if user.role else ""
    if role not in ("admin", "ops_manager"):
        raise HTTPException(403, "Admin or Ops Manager access required")

def any_auth(user):
    if not user:
        raise HTTPException(401, "Authentication required")

# ─── Templates ───

@router.get("/templates")
async def list_templates(db: AsyncSession = Depends(get_db), user: User = Depends(get_current_user)):
    any_auth(user)
    result = await db.execute(
        select(ReportTemplate).options(selectinload(ReportTemplate.creator)).order_by(ReportTemplate.created_at.desc())
    )
    templates = result.scalars().all()
    return [{
        "id": t.id,
        "name": t.name,
        "description": t.description,
        "period_type": t.period_type,
        "sections": json.loads(t.sections) if isinstance(t.sections, str) else t.sections,
        "created_by_name": t.creator.name if t.creator else None,
        "created_at": str(t.created_at) if t.created_at else None,
    } for t in templates]

@router.post("/templates/create")
async def create_template(
    name: str = Form(...),
    description: str = Form(None),
    period_type: str = Form("monthly"),
    sections: str = Form(...),
    db: AsyncSession = Depends(get_db),
    user: User = Depends(get_current_user),
):
    admin_or_ops(user)
    if not name or not sections:
        raise HTTPException(400, "Name and sections required")
    try:
        json.loads(sections)
    except:
        raise HTTPException(400, "Sections must be valid JSON")
    t = ReportTemplate(
        name=name,
        description=description,
        period_type=period_type,
        sections=sections,
        created_by=user.id,
    )
    db.add(t)
    await db.commit()
    await db.refresh(t)
    return {"id": t.id, "name": t.name, "msg": "Template created"}

@router.put("/templates/{template_id}")
async def update_template(
    template_id: int,
    name: str = Form(None),
    description: str = Form(None),
    period_type: str = Form(None),
    sections: str = Form(None),
    db: AsyncSession = Depends(get_db),
    user: User = Depends(get_current_user),
):
    admin_or_ops(user)
    result = await db.execute(select(ReportTemplate).where(ReportTemplate.id == template_id))
    t = result.scalar_one_or_none()
    if not t:
        raise HTTPException(404, "Template not found")
    if name: t.name = name
    if description is not None: t.description = description
    if period_type: t.period_type = period_type
    if sections:
        try:
            json.loads(sections)
        except:
            raise HTTPException(400, "Sections must be valid JSON")
        t.sections = sections
    await db.commit()
    return {"msg": "Template updated"}

@router.delete("/templates/{template_id}")
async def delete_template(template_id: int, db: AsyncSession = Depends(get_db), user: User = Depends(get_current_user)):
    admin_or_ops(user)
    result = await db.execute(select(ReportTemplate).where(ReportTemplate.id == template_id))
    t = result.scalar_one_or_none()
    if not t:
        raise HTTPException(404, "Template not found")
    await db.delete(t)
    await db.commit()
    return {"msg": "Template deleted"}

# ─── Periods ───

@router.get("/periods")
async def list_periods(template_id: int = None, db: AsyncSession = Depends(get_db), user: User = Depends(get_current_user)):
    any_auth(user)
    q = select(ReportPeriod).order_by(ReportPeriod.created_at.desc())
    if template_id:
        q = q.where(ReportPeriod.template_id == template_id)
    result = await db.execute(q)
    periods = result.scalars().all()
    return [{
        "id": p.id,
        "template_id": p.template_id,
        "label": p.label,
        "start_date": str(p.start_date) if p.start_date else None,
        "end_date": str(p.end_date) if p.end_date else None,
        "is_open": p.is_open,
        "section_assignments": json.loads(p.section_assignments) if p.section_assignments else {},
        "created_at": str(p.created_at) if p.created_at else None,
    } for p in periods]

@router.post("/periods/create")
async def create_period(
    template_id: int = Form(...),
    label: str = Form(...),
    start_date: str = Form(None),
    end_date: str = Form(None),
    db: AsyncSession = Depends(get_db),
    user: User = Depends(get_current_user),
):
    admin_or_ops(user)
    p = ReportPeriod(
        template_id=template_id,
        label=label,
        start_date=datetime.strptime(start_date, "%Y-%m-%d").date() if start_date else None,
        end_date=datetime.strptime(end_date, "%Y-%m-%d").date() if end_date else None,
    )
    db.add(p)
    await db.commit()
    await db.refresh(p)
    return {"id": p.id, "label": p.label, "msg": "Period created"}

@router.post("/periods/{period_id}/close")
async def close_period(period_id: int, db: AsyncSession = Depends(get_db), user: User = Depends(get_current_user)):
    admin_or_ops(user)
    result = await db.execute(select(ReportPeriod).where(ReportPeriod.id == period_id))
    p = result.scalar_one_or_none()
    if not p:
        raise HTTPException(404, "Period not found")
    p.is_open = False
    await db.commit()
    return {"msg": "Period closed"}

@router.put("/periods/{period_id}/assignments")
async def update_assignments(
    period_id: int,
    assignments: str = Form(...),
    db: AsyncSession = Depends(get_db),
    user: User = Depends(get_current_user),
):
    admin_or_ops(user)
    try:
        parsed = json.loads(assignments)
    except:
        raise HTTPException(400, "Invalid JSON")
    result = await db.execute(select(ReportPeriod).where(ReportPeriod.id == period_id))
    p = result.scalar_one_or_none()
    if not p:
        raise HTTPException(404, "Period not found")
    p.section_assignments = json.dumps(parsed)
    await db.commit()
    return {"msg": "Assignments updated"}

@router.get("/periods/{period_id}/assignments")
async def get_assignments(period_id: int, db: AsyncSession = Depends(get_db), user: User = Depends(get_current_user)):
    any_auth(user)
    result = await db.execute(select(ReportPeriod).where(ReportPeriod.id == period_id))
    p = result.scalar_one_or_none()
    if not p:
        raise HTTPException(404, "Period not found")
    return json.loads(p.section_assignments) if p.section_assignments else {}

# ─── Submissions ───

@router.get("/submissions")
async def list_submissions(
    period_id: int = None,
    section_key: str = None,
    assigned_to: int = None,
    db: AsyncSession = Depends(get_db),
    user: User = Depends(get_current_user),
):
    any_auth(user)
    q = select(ReportSubmission).options(
        selectinload(ReportSubmission.period),
        selectinload(ReportSubmission.assignee),
    ).order_by(ReportSubmission.created_at.desc())
    if period_id: q = q.where(ReportSubmission.period_id == period_id)
    if section_key: q = q.where(ReportSubmission.section_key == section_key)
    if assigned_to: q = q.where(ReportSubmission.assigned_to == assigned_to)
    result = await db.execute(q)
    subs = result.scalars().all()
    return [{
        "id": s.id,
        "period_id": s.period_id,
        "section_key": s.section_key,
        "assigned_to": s.assigned_to,
        "assigned_to_name": s.assignee.name if s.assignee else None,
        "field_values": json.loads(s.field_values) if isinstance(s.field_values, str) else s.field_values,
        "status": s.status,
        "submitted_at": str(s.submitted_at) if s.submitted_at else None,
        "created_at": str(s.created_at) if s.created_at else None,
    } for s in subs]

@router.post("/submissions/save")
async def save_submission(
    period_id: int = Form(...),
    section_key: str = Form(...),
    assigned_to: int = Form(None),
    field_values: str = Form(...),
    status: str = Form("draft"),
    db: AsyncSession = Depends(get_db),
    user: User = Depends(get_current_user),
):
    any_auth(user)
    if status not in ("draft", "submitted"):
        raise HTTPException(400, "Invalid status")
    try:
        fv = json.loads(field_values)
    except:
        raise HTTPException(400, "field_values must be valid JSON")
    result = await db.execute(
        select(ReportSubmission).where(
            ReportSubmission.period_id == period_id,
            ReportSubmission.section_key == section_key,
        )
    )
    existing = result.scalar_one_or_none()
    if existing:
        existing.field_values = field_values
        existing.status = status
        existing.assigned_to = assigned_to or existing.assigned_to
        if status == "submitted":
            existing.submitted_at = datetime.utcnow()
    else:
        sub = ReportSubmission(
            period_id=period_id,
            section_key=section_key,
            assigned_to=assigned_to,
            field_values=field_values,
            status=status,
            submitted_at=datetime.utcnow() if status == "submitted" else None,
        )
        db.add(sub)
    await db.commit()
    return {"msg": "Submission saved", "status": status}

# ─── Export ───

@router.get("/export/{period_id}")
async def export_report(
    period_id: int,
    format: str = "json",
    db: AsyncSession = Depends(get_db),
    user: User = Depends(get_current_user),
):
    any_auth(user)
    result = await db.execute(
        select(ReportPeriod).options(selectinload(ReportPeriod.template)).where(ReportPeriod.id == period_id)
    )
    period = result.scalar_one_or_none()
    if not period:
        raise HTTPException(404, "Period not found")

    template_sections = json.loads(period.template.sections) if isinstance(period.template.sections, str) else period.template.sections

    result = await db.execute(
        select(ReportSubmission).options(selectinload(ReportSubmission.assignee)).where(ReportSubmission.period_id == period_id)
    )
    subs = result.scalars().all()

    submission_map = {}
    for s in subs:
        submission_map[s.section_key] = {
            "field_values": json.loads(s.field_values) if isinstance(s.field_values, str) else s.field_values,
            "status": s.status,
            "assigned_to_name": s.assignee.name if s.assignee else "Unassigned",
            "submitted_at": str(s.submitted_at) if s.submitted_at else "—",
        }

    period_assignments = json.loads(period.section_assignments) if period.section_assignments else {}

    # Build user-id → name mapping
    result = await db.execute(select(User))
    all_users = result.scalars().all()
    user_names = {u.id: u.name for u in all_users}

    report_data = {
        "template_name": period.template.name,
        "period_label": period.label,
        "period_start": str(period.start_date) if period.start_date else None,
        "period_end": str(period.end_date) if period.end_date else None,
        "sections": [],
    }
    for sec in template_sections:
        sec_key = sec.get("key", "")
        sub_data = submission_map.get(sec_key, {})
        assigned_uid = period_assignments.get(sec_key) or sub_data.get("assigned_to")
        assigned_name = user_names.get(assigned_uid) if assigned_uid else sub_data.get("assigned_to_name", "Unassigned")
        report_data["sections"].append({
            "title": sec.get("title", sec_key),
            "key": sec_key,
            "fields": sec.get("fields", []),
            "values": sub_data.get("field_values", {}),
            "status": sub_data.get("status", "pending"),
            "assigned_to": assigned_name,
            "submitted_at": sub_data.get("submitted_at", "—"),
        })

    if format == "json":
        return report_data

    elif format == "text":
        lines = []
        lines.append(f"{'='*60}")
        lines.append(f"  {report_data['template_name']}")
        lines.append(f"  Period: {report_data['period_label']}")
        lines.append(f"{'='*60}\n")
        for sec in report_data["sections"]:
            lines.append(f"── {sec['title']} ──")
            lines.append(f"   Assigned to: {sec['assigned_to']}")
            lines.append(f"   Status: {sec['status']}")
            for f in sec["fields"]:
                fk = f.get("key", "")
                fl = f.get("label", fk)
                val = sec["values"].get(fk, "—")
                lines.append(f"   {fl}: {val}")
            lines.append("")
        return Response("\n".join(lines), media_type="text/plain")

    elif format == "html":
        rows = ""
        for sec in report_data["sections"]:
            fields_html = ""
            for f in sec["fields"]:
                fk = f.get("key", "")
                fl = f.get("label", fk)
                val = sec["values"].get(fk, "—")
                fields_html += f"<tr><td style='padding:4px 8px;color:#555;font-weight:600'>{fl}</td><td style='padding:4px 8px'>{val}</td></tr>"
            status_color = {"draft":"#E65100","submitted":"#1565C0","pending":"#999"}
            rows += f"""
            <div style="margin-bottom:16px;border:1px solid #e0e4e8;border-radius:8px;padding:12px">
                <div style="font-size:14px;font-weight:700;color:#0b3d91;margin-bottom:4px">{sec['title']}</div>
                <div style="font-size:12px;color:#888;margin-bottom:8px">Assigned: {sec['assigned_to']} | Status: <span style="color:{status_color.get(sec['status'],'#999')}">{sec['status']}</span></div>
                <table style="width:100%;border-collapse:collapse;font-size:13px">{fields_html}</table>
            </div>"""
        html = f"""<html><head><meta charset="utf-8"><title>{report_data['template_name']}</title></head><body style="font-family:sans-serif;padding:20px;max-width:800px;margin:auto">
            <h1 style="color:#0b3d91">{report_data['template_name']}</h1>
            <p style="color:#666">Period: {report_data['period_label']}</p>
            {rows}
        </body></html>"""
        return Response(html, media_type="text/html")

    elif format == "docx":
        try:
            from docx import Document
            from docx.shared import Inches, Pt
            doc = Document()
            doc.add_heading(report_data["template_name"], level=1)
            doc.add_paragraph(f"Period: {report_data['period_label']}")
            for sec in report_data["sections"]:
                doc.add_heading(sec["title"], level=2)
                doc.add_paragraph(f"Assigned to: {sec['assigned_to']}  |  Status: {sec['status']}")
                table = doc.add_table(rows=1, cols=2)
                table.style = "Light Grid Accent 1"
                hdr = table.rows[0].cells
                hdr[0].text = "Field"
                hdr[1].text = "Value"
                for f in sec["fields"]:
                    fk = f.get("key", "")
                    fl = f.get("label", fk)
                    val = sec["values"].get(fk, "—")
                    row = table.add_row().cells
                    row[0].text = fl
                    row[1].text = str(val)
                doc.add_paragraph("")
            import io
            buf = io.BytesIO()
            doc.save(buf)
            buf.seek(0)
            return Response(
                buf.read(),
                media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                headers={"Content-Disposition": f"attachment; filename={period.template.name.replace(' ','_')}_{period.label.replace(' ','_')}.docx"}
            )
        except ImportError:
            raise HTTPException(500, "DOCX export not available (python-docx not installed)")

    elif format == "pptx":
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = report_data["template_name"]
            slide.placeholders[1].text = f"Period: {report_data['period_label']}"
            for sec in report_data["sections"]:
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                slide.shapes.title.text = sec["title"]
                txBox = slide.placeholders[1]
                tf = txBox.text_frame
                tf.text = f"Assigned to: {sec['assigned_to']}\nStatus: {sec['status']}\n"
                for f in sec["fields"]:
                    fk = f.get("key", "")
                    fl = f.get("label", fk)
                    val = sec["values"].get(fk, "—")
                    p = tf.add_paragraph()
                    p.text = f"{fl}: {val}"
                    p.space_after = Pt(4)
            import io
            buf = io.BytesIO()
            prs.save(buf)
            buf.seek(0)
            return Response(
                buf.read(),
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                headers={"Content-Disposition": f"attachment; filename={period.template.name.replace(' ','_')}_{period.label.replace(' ','_')}.pptx"}
            )
        except ImportError:
            raise HTTPException(500, "PPTX export not available (python-pptx not installed)")

    elif format == "pdf":
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.lib import colors
            import io
            buf = io.BytesIO()
            doc = SimpleDocTemplate(buf, pagesize=A4)
            styles = getSampleStyleSheet()
            elements = []
            elements.append(Paragraph(report_data["template_name"], styles["Title"]))
            elements.append(Paragraph(f"Period: {report_data['period_label']}", styles["Normal"]))
            elements.append(Spacer(1, 12))
            for sec in report_data["sections"]:
                elements.append(Paragraph(sec["title"], styles["Heading2"]))
                elements.append(Paragraph(f"Assigned to: {sec['assigned_to']} | Status: {sec['status']}", styles["Normal"]))
                elements.append(Spacer(1, 6))
                data = [["Field", "Value"]]
                for f in sec["fields"]:
                    fk = f.get("key", "")
                    fl = f.get("label", fk)
                    val = sec["values"].get(fk, "—")
                    data.append([fl, str(val)])
                tbl = Table(data, colWidths=[200, 350])
                tbl.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#0b3d91")),
                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("FONTSIZE", (0, 0), (-1, -1), 10),
                    ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                    ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f5f7fa")]),
                ]))
                elements.append(tbl)
                elements.append(Spacer(1, 12))
            doc.build(elements)
            buf.seek(0)
            return Response(
                buf.read(),
                media_type="application/pdf",
                headers={"Content-Disposition": f"attachment; filename={period.template.name.replace(' ','_')}_{period.label.replace(' ','_')}.pdf"}
            )
        except ImportError:
            raise HTTPException(500, "PDF export not available (reportlab not installed)")

    raise HTTPException(400, f"Unsupported format: {format}")
