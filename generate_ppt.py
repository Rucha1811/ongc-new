from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BLUE = RGBColor(0x0B, 0x3D, 0x91)
MED_BLUE = RGBColor(0x15, 0x65, 0xC0)
LIGHT_BLUE = RGBColor(0xE3, 0xEE, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x5A, 0x6A, 0x7A)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
DARK_GRAY = RGBColor(0x34, 0x40, 0x54)
GREEN = RGBColor(0x1B, 0x5E, 0x20)
RED = RGBColor(0xB7, 0x1C, 0x1C)
ORANGE = RGBColor(0xE6, 0x51, 0x00)
PURPLE = RGBColor(0x7B, 0x1F, 0xA2)

def add_bg(slide, color=LIGHT_GRAY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header_bar(slide, title, subtitle=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(10), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(0.75), Inches(10), Inches(0.4))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
        p2.font.italic = True

def add_screenshot_placeholder(slide, left, top, width, height, label):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xED, 0xF2)
    shape.line.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
    shape.line.width = Pt(1.5)
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = f"[ {label} ]"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.font.italic = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

def add_bullet_text(slide, left, top, width, height, items, font_size=16, color=BLACK, spacing=Pt(6), bold_prefix=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(item, tuple):
            # (bold_part, normal_part)
            if bold_prefix:
                run = p.add_run()
                run.text = item[0]
                run.font.bold = True
                run.font.size = Pt(font_size)
                run.font.color.rgb = color
                run = p.add_run()
                run.text = item[1]
                run.font.size = Pt(font_size)
                run.font.color.rgb = color
            else:
                p.text = item[0] + item[1]
                p.font.size = Pt(font_size)
                p.font.color.rgb = color
        else:
            p.text = f"  {item}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
        p.space_after = spacing
        p.level = 0
    return txBox

def add_footer(slide, text="Advance Data Repository by Geophysical Services"):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.font.italic = True

def add_table(slide, left, top, width, height, headers, rows, col_widths=None):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    # Header
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
    # Rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BLUE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = BLACK
                p.alignment = PP_ALIGN.CENTER
    return table_shape

# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BLUE)

txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.3), Inches(2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Advance Data Repository"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "by Geophysical Services"
p2.font.size = Pt(28)
p2.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(8)

p3 = tf.add_paragraph()
p3.text = "Digital Platform for Secure Storage, Data Management & Access"
p3.font.size = Pt(18)
p3.font.color.rgb = RGBColor(0x90, 0xCA, 0xF9)
p3.alignment = PP_ALIGN.CENTER
p3.space_before = Pt(24)

p4 = tf.add_paragraph()
p4.text = "User Manual"
p4.font.size = Pt(22)
p4.font.bold = True
p4.font.color.rgb = WHITE
p4.alignment = PP_ALIGN.CENTER
p4.space_before = Pt(40)

# ============================================================
# SLIDE 2: TABLE OF CONTENTS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Table of Contents")
items = [
    "System Overview & Problem Statement",
    "System Architecture",
    "User Roles & Access Control",
    "Authentication & Login",
    "Dashboard",
    "File Upload",
    "File Records & Search",
    "Approval Workflow",
    "Notifications",
    "User Management (Admin)",
    "Settings & Dropdown Manager (Admin)",
    "Reports",
    "Security & Design Highlights",
    "Technology Stack",
    "Demo Flow",
]
add_bullet_text(slide, Inches(1), Inches(1.6), Inches(11.3), Inches(5.5), items, font_size=20, color=DARK_BLUE, spacing=Pt(10))
add_footer(slide)

# ============================================================
# SLIDE 3: OVERVIEW
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "System Overview", "Problem Statement & Solution")

left_content = [
    ("Problem: ", "Geophysical Services deals with massive volumes of exploration data — seismic, well data, reports, contracts, maps. Data was scattered across shared folders with no centralized access control, audit trail, or approval workflow."),
    ("Solution: ", "A web-based Advance Data Repository with role-based access, file upload/approval workflow, notifications, and centralized management."),
    ("Users: ", "4 roles — Admin, Operations Manager, Data Creator, Viewer — each with distinct permissions and UI."),
]
add_bullet_text(slide, Inches(0.6), Inches(1.5), Inches(7.5), Inches(4.5), left_content, font_size=14, color=BLACK, spacing=Pt(14))
add_screenshot_placeholder(slide, Inches(8.5), Inches(1.5), Inches(4.3), Inches(4.5), "Screenshot: Login Page")
add_footer(slide)

# ============================================================
# SLIDE 4: SYSTEM ARCHITECTURE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "System Architecture")
items = [
    ("Frontend: ", "React 18 + Vite — runs on localhost:5173"),
    ("Backend: ", "Python FastAPI — REST API on localhost:8000"),
    ("Database: ", "PostgreSQL 15 + pgvector (Docker, port 5433)"),
    ("File Storage: ", "Disk (uploads/{category}/{classification}/{filename}) + DB (BYTEA fallback)"),
    ("Admin Tool: ", "pgAdmin 4 on localhost:5050"),
    ("OCR Engine: ", "PP-OCRv5 (Python 3.13) for scanned documents — optional"),
]
add_bullet_text(slide, Inches(0.6), Inches(1.5), Inches(7.5), Inches(5.0), items, font_size=15, color=BLACK, spacing=Pt(12))
add_screenshot_placeholder(slide, Inches(8.5), Inches(1.5), Inches(4.3), Inches(3.0), "Screenshot: pgAdmin tables")
add_screenshot_placeholder(slide, Inches(8.5), Inches(4.8), Inches(4.3), Inches(2.0), "Screenshot: Folder structure")
add_footer(slide)

# ============================================================
# SLIDE 5: USER ROLES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "User Roles & Access Control")
add_table(slide, Inches(0.6), Inches(1.6), Inches(12.1), Inches(3.5),
    ["Role", "Classification Access", "Permissions"],
    [
        ["Admin\n(Full Control)", "General / Available for All\nSensitive / Internal Use\nConfidential\nHighly Confidential / Restricted", "Upload, approve, reject, manage users,\nmanage settings, grant permissions"],
        ["Operations Manager", "General / Available for All\nSensitive / Internal Use", "Upload, approve, reject files"],
        ["Data Creator", "General / Available for All\n+ own uploads (any classification)", "Upload files, view own files\n& General files"],
        ["End User / Viewer", "General / Available for All", "View & download General files\n(grants can extend access)"],
    ],
    col_widths=[Inches(2.2), Inches(4.5), Inches(5.4)]
)
note_items = [
    "Role IDs are fixed: 1=admin, 2=ops_manager, 3=data_creator, 4=viewer",
    "Admin role (name='admin') cannot be changed; all other users can be promoted to admin",
    "Temporary permission grants expire after 1 hour",
]
add_bullet_text(slide, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.5), note_items, font_size=13, color=GRAY, spacing=Pt(4))
add_footer(slide)

# ============================================================
# SLIDE 6: AUTHENTICATION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Authentication & Login")
items = [
    ("Login: ", "CPF (unique employee ID) + password"),
    ("Pre-seeded Users:", ""),
]
add_bullet_text(slide, Inches(0.6), Inches(1.5), Inches(7.5), Inches(1.0), items, font_size=15, color=BLACK, spacing=Pt(4))

add_table(slide, Inches(0.6), Inches(2.5), Inches(7.5), Inches(2.8),
    ["CPF", "Name", "Role", "Password"],
    [
        ["100001", "Sh. Sandip Kumar Kaur", "Admin", "admin123"],
        ["100002", "Sh. Pankaj Kumar", "Ops Manager", "ops123"],
        ["100003", "Sh. Manjunath B. V.", "Data Creator", "user123"],
        ["100004", "Ms. T. Sushma", "Viewer", "view123"],
        ["100005", "Rucha", "Admin", "Rucha"],
    ],
    col_widths=[Inches(1.5), Inches(2.4), Inches(1.7), Inches(1.9)]
)
items2 = [
    "Password hashed with bcrypt ($2b$12$ rounds)",
    "Session persists via JWT token in sessionStorage",
]
add_bullet_text(slide, Inches(0.6), Inches(5.6), Inches(7.5), Inches(1.0), items2, font_size=13, color=GRAY, spacing=Pt(2))
add_screenshot_placeholder(slide, Inches(8.5), Inches(1.5), Inches(4.3), Inches(4.5), "Screenshot: Login Page")
add_footer(slide)

# ============================================================
# SLIDE 7: DASHBOARD
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Dashboard", "Role-specific landing page")
items = [
    ("Admin: ", "Total files, pending approvals, recent uploads, bar charts"),
    ("Ops Manager: ", "Files pending approval, recent uploads summary"),
    ("Data Creator: ", "Own files summary, recent uploads"),
    ("Viewer: ", "Available files summary"),
    "Stat cards with gradient backgrounds, consistent font & spacing",
]
add_bullet_text(slide, Inches(0.6), Inches(1.5), Inches(7.5), Inches(3.5), items, font_size=15, color=BLACK, spacing=Pt(8))
add_screenshot_placeholder(slide, Inches(8.5), Inches(1.5), Inches(4.3), Inches(2.5), "Screenshot: Admin Dashboard")
add_screenshot_placeholder(slide, Inches(8.5), Inches(4.3), Inches(4.3), Inches(2.5), "Screenshot: Creator Dashboard")
add_footer(slide)

# ============================================================
# SLIDE 8: FILE UPLOAD
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "File Upload", "Structured metadata for every file")
items = [
    "Upload form with 7 dropdown fields:",
    ("  Section: ", "GP-03, GP-06, REL, RCC, HSE, …"),
    ("  Category: ", "53 categories — General Admin, Accounts, Seismic Data, …"),
    ("  Season: ", "2025-26 through 1956-57"),
    ("  Block: ", "Ankleshwar, Ahmedabad, Mehsana, Rajasthan, Other"),
    ("  File Type: ", "PDF, DOCX, XLSX, CSV, ZIP, …"),
    ("  Data Type: ", "Seismic 2D/3D, LFPS, VSP, Any Other Data"),
    ("  Classification: ", "4 levels — General to Highly Confidential"),
    "",
    "All dropdown values stored in central 'lookups' DB table",
    "File saved to disk: uploads/{category_sanitized}/{classification_sanitized}/{filename}",
    "File data also stored in DB as BYTEA for redundancy",
]
add_bullet_text(slide, Inches(0.6), Inches(1.5), Inches(7.5), Inches(5.5), items, font_size=13, color=BLACK, spacing=Pt(3))
add_screenshot_placeholder(slide, Inches(8.5), Inches(1.5), Inches(4.3), Inches(5.0), "Screenshot: Upload File form")
add_footer(slide)

# ============================================================
# SLIDE 9: APPROVAL WORKFLOW
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Approval Workflow", "Data Creator → Ops Manager / Admin")
items = [
    "1. Data Creator uploads file → status = Pending",
    "2. Ops Manager / Admin reviews in 'Pending Approval' page",
    "3. Two actions:",
    ("   Approve: ", "Optionally change classification at approval time"),
    ("   Reject: ", "Must provide a reason (comment required)"),
    "4. Uploader gets a notification in both cases",
    "5. Approved files appear under 'Approved Files' for all authorized roles",
    "",
    "Classification color-coding in table:",
    "  General = Green   |   Sensitive = Orange   |   Confidential = Red   |   Highly Confidential = Purple",
]
add_bullet_text(slide, Inches(0.6), Inches(1.5), Inches(7.5), Inches(5.0), items, font_size=14, color=BLACK, spacing=Pt(5))
add_screenshot_placeholder(slide, Inches(8.5), Inches(1.5), Inches(4.3), Inches(2.2), "Screenshot: Approve dialog")
add_screenshot_placeholder(slide, Inches(8.5), Inches(4.0), Inches(4.3), Inches(2.2), "Screenshot: Reject dialog")
add_footer(slide)

# ============================================================
# SLIDE 10: NOTIFICATIONS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Notifications", "Real-time alerts for file status changes")
items = [
    "Bell icon in header showing unread notification count",
    "Click bell → dropdown with all notifications",
    "Each notification shows: file name, action taken, who did it, reason (if rejected)",
    "Mark individual notification as read, or 'Mark all read'",
    "Auto-polls every 15 seconds for new notifications",
    "Notifications created for:",
    ("  Approval: ", "\"Your file X has been approved by Y\""),
    ("  Rejection: ", "\"Your file X was rejected by Y. Reason: …\""),
]
add_bullet_text(slide, Inches(0.6), Inches(1.5), Inches(7.5), Inches(4.5), items, font_size=14, color=BLACK, spacing=Pt(6))
add_screenshot_placeholder(slide, Inches(8.5), Inches(1.5), Inches(4.3), Inches(3.0), "Screenshot: Notifications dropdown")
add_screenshot_placeholder(slide, Inches(8.5), Inches(4.8), Inches(4.3), Inches(1.5), "Screenshot: Bell icon with count badge")
add_footer(slide)

# ============================================================
# SLIDE 11: FILE RECORDS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "File Records & Search", "Browse, filter, and download files")
items = [
    "Table view with columns: S.No, File Name, Category, Classification,\n  Section, Season, Uploaded By, Status, Actions",
    "Status filters: All, Pending, Approved, Rejected",
    "Search by file name",
    "Classification color-coded badges",
    "Download button for approved files",
    "Click file name to view details",
    "Sidebar navigation:",
    ("  Admin: ", "Dashboard | Upload File | File Records | Pending Approval | Approved Files | Rejected Files | Reports | Users | Access Permissions | Settings"),
    ("  Ops Manager: ", "Dashboard | Upload File | File Records | Pending Approval | Approved Files | Rejected Files | Reports"),
    ("  Data Creator: ", "Dashboard | Upload File | My Files | Reports"),
    ("  Viewer: ", "Dashboard | File Records | Approved Files | Reports"),
]
add_bullet_text(slide, Inches(0.6), Inches(1.5), Inches(7.5), Inches(5.5), items, font_size=12, color=BLACK, spacing=Pt(3))
add_screenshot_placeholder(slide, Inches(8.5), Inches(1.5), Inches(4.3), Inches(3.5), "Screenshot: File Records table")
add_screenshot_placeholder(slide, Inches(8.5), Inches(5.3), Inches(4.3), Inches(1.3), "Screenshot: Sidebar menus")
add_footer(slide)

# ============================================================
# SLIDE 12: USER MANAGEMENT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "User Management", "Admin only — role assignment & permissions")
items = [
    "View all users in table format",
    "Change user role via dropdown (only 'admin' role name is immutable)",
    "Any user can be promoted to admin (no restriction)",
    "Grant temporary classification access:",
    ("  Expires after: ", "1 hour"),
    ("  Filtered automatically: ", "Expired permissions excluded from API responses"),
    "Activity log tracks user changes",
]
add_bullet_text(slide, Inches(0.6), Inches(1.5), Inches(7.5), Inches(3.5), items, font_size=14, color=BLACK, spacing=Pt(8))
add_screenshot_placeholder(slide, Inches(8.5), Inches(1.5), Inches(4.3), Inches(2.5), "Screenshot: User Management page")
add_screenshot_placeholder(slide, Inches(8.5), Inches(4.3), Inches(4.3), Inches(2.5), "Screenshot: Grant permission dialog")
add_footer(slide)

# ============================================================
# SLIDE 13: SETTINGS — DROPDOWN MANAGER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Settings: Dropdown Options Manager", "Admin only — manage dropdown values dynamically")
items = [
    "All dropdown options stored in central 'lookups' DB table",
    "7 option types: Sections, Categories, Seasons, Blocks, File Types, Data Types, Classifications",
    "Admin can Add / Edit / Delete values for any type",
    "Changes reflect immediately in Upload forms",
    "156 pre-seeded values covering all dropdowns",
    "No hardcoded arrays — fully data-driven",
]
add_bullet_text(slide, Inches(0.6), Inches(1.5), Inches(7.5), Inches(3.5), items, font_size=14, color=BLACK, spacing=Pt(8))
add_screenshot_placeholder(slide, Inches(8.5), Inches(1.5), Inches(4.3), Inches(4.0), "Screenshot: Settings / Dropdown Manager")
add_footer(slide)

# ============================================================
# SLIDE 14: REPORTS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Reports", "Monthly file statistics")
items = [
    "Monthly file upload statistics with bar charts",
    "Breakdown by month (last 12 months)",
    "Available to all roles",
    "Data aggregated from files table by created_at month",
]
add_bullet_text(slide, Inches(0.6), Inches(1.5), Inches(7.5), Inches(3.0), items, font_size=15, color=BLACK, spacing=Pt(8))
add_screenshot_placeholder(slide, Inches(1.5), Inches(3.5), Inches(10.3), Inches(3.0), "Screenshot: Reports page with bar chart")
add_footer(slide)

# ============================================================
# SLIDE 15: SECURITY & DESIGN
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Security & Design Highlights")
items = [
    ("Integer-only IDs: ", "All primary/foreign keys are clean auto-increment integers — no UUIDs"),
    ("Role-Based Access: ", "Each role sees only permitted classifications"),
    ("Expiring Permissions: ", "Admin-granted permissions auto-expire after 1 hour"),
    ("Admin Role Immutable: ", "The role named 'admin' cannot be changed/deleted"),
    ("Rejection Requires Reason: ", "Mandatory comment field prevents silent rejections"),
    ("Password Security: ", "bcrypt hashing ($2b$12$ rounds), JWT token auth"),
    ("Professional UI: ", "Segoe UI font, gradient headers, card-based layout, consistent spacing"),
    ("Footer Hidden: ", "Footer removed after login for cleaner workspace"),
    ("Font Sizes: ", "Table text 14px, small elements 13px, tiny 12px"),
]
add_bullet_text(slide, Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.5), items, font_size=15, color=BLACK, spacing=Pt(8))
add_footer(slide)

# ============================================================
# SLIDE 16: TECHNOLOGY STACK
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Technology Stack")
add_table(slide, Inches(0.6), Inches(1.6), Inches(12.1), Inches(4.0),
    ["Layer", "Technology", "Version", "Port"],
    [
        ["Frontend", "React + Vite", "React 18", "5173"],
        ["Backend API", "Python FastAPI", "Python 3.14", "8000"],
        ["Database", "PostgreSQL + pgvector", "PostgreSQL 15", "5433 (Docker)"],
        ["Admin Tool", "pgAdmin 4", "Latest", "5050"],
        ["OCR (optional)", "PP-OCRv5 (PaddleOCR)", "Python 3.13", "—"],
        ["Embeddings (optional)", "sentence-transformers", "all-MiniLM-L6-v2", "—"],
        ["Container", "Docker Compose", "—", "—"],
    ],
    col_widths=[Inches(2.5), Inches(4.5), Inches(2.5), Inches(2.6)]
)
items2 = [
    "Backend: uvicorn ASGI server, SQLAlchemy async ORM, Alembic migrations",
    "Database: pgvector extension for vector embeddings (optional semantic search)",
    "Frontend: React hooks (useState, useEffect, useCallback), inline styles via JS object",
]
add_bullet_text(slide, Inches(0.6), Inches(5.9), Inches(12.1), Inches(1.2), items2, font_size=13, color=GRAY, spacing=Pt(3))
add_footer(slide)

# ============================================================
# SLIDE 17: DEMO FLOW
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, "Live Demo Walkthrough")
items = [
    ("1. Login as Admin:", " CPF: 100001 / Password: admin123"),
    "   → Show Dashboard stats, sidebar menu items",
    "   → Go to User Management → show user table + role dropdown",
    "   → Go to Settings → show Dropdown Options Manager",
    "",
    ("2. Login as Data Creator:", " CPF: 100003 / Password: user123"),
    "   → Upload a new file (fill all dropdowns, choose 'General')",
    "   → Go to My Files → confirm file shows as Pending",
    "",
    ("3. Login as Ops Manager:", " CPF: 100002 / Password: ops123"),
    "   → Go to Pending Approval → find the file → Approve it",
    "   → (or Reject with a reason)",
    "",
    ("4. Login back as Data Creator:", ""),
    "   → See notification bell with red badge",
    "   → Click bell → see approval/rejection notification",
    "",
    ("5. Login as Viewer:", " CPF: 100004 / Password: view123"),
    "   → Browse Approved Files → download the file",
]
add_bullet_text(slide, Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.8), items, font_size=14, color=BLACK, spacing=Pt(3))
add_footer(slide)

# ============================================================
# SLIDE 18: THANK YOU
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)

txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Thank You"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "Advance Data Repository by Geophysical Services"
p2.font.size = Pt(22)
p2.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(20)

p3 = tf.add_paragraph()
p3.text = "Questions?"
p3.font.size = Pt(20)
p3.font.color.rgb = WHITE
p3.alignment = PP_ALIGN.CENTER
p3.space_before = Pt(30)

# SAVE
output_path = "/Users/ruchatejaskumargandhi/Desktop/ONGC 3/ONGC_Advance_Data_Repository_Manual.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
